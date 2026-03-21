#!/usr/bin/env python3
"""Generate PPTX for Финансовый компас 1.0 — in FinDirector template style"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os
import copy
from lxml import etree

base = "/Users/yerkink/Karina-Media/prezentacia/"

# Use template for style
prs = Presentation(base + "template.pptx")

# Remove all existing slides
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[-1].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    if rId is None:
        rId_attr = list(prs.slides._sldIdLst[-1].attrib.keys())
        for attr in rId_attr:
            if 'id' in attr.lower() and 'r' in attr.lower():
                rId = prs.slides._sldIdLst[-1].get(attr)
                break
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[-1])

# Colors matching template
BG = RGBColor(0x03, 0x0A, 0x18)
BLUE = RGBColor(0x54, 0x6E, 0xE5)
WHITE = RGBColor(0xF5, 0xF6, 0xFA)
GRAY = RGBColor(0x7A, 0x8C, 0xB3)
GREEN = RGBColor(0x10, 0xB9, 0x81)
RED = RGBColor(0xEF, 0x44, 0x44)
GOLD = RGBColor(0xF5, 0x9E, 0x0B)
FONT = 'Calibri'

layout = prs.slide_layouts[0]

def add_slide():
    slide = prs.slides.add_slide(layout)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    return slide

def add_title(slide, text, size=26):
    """Blue bold title at top"""
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9.0), Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = BLUE
    run.font.name = FONT

def add_body(slide, lines, left=0.5, top=1.3, width=9.0, height=3.6, text_size=14, bullet_style=True):
    """Body text with blue bullet points"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True

    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        p.space_before = Pt(4)
        p.space_after = Pt(2)

        if isinstance(line, tuple):
            # (text, size, bold, color)
            text, sz, bold, color = line
            run = p.add_run()
            run.text = text
            run.font.size = Pt(sz)
            run.font.bold = bold
            run.font.color.rgb = color
            run.font.name = FONT
        elif line.startswith('•'):
            # Bullet point: blue bullet + white text
            run_bullet = p.add_run()
            run_bullet.text = "• "
            run_bullet.font.size = Pt(text_size)
            run_bullet.font.bold = True
            run_bullet.font.color.rgb = BLUE
            run_bullet.font.name = FONT

            run_text = p.add_run()
            run_text.text = line[2:] if line.startswith('• ') else line[1:]
            run_text.font.size = Pt(text_size)
            run_text.font.bold = False
            run_text.font.color.rgb = WHITE
            run_text.font.name = FONT
        elif line == '':
            pass  # empty line
        else:
            run = p.add_run()
            run.text = line
            run.font.size = Pt(text_size)
            run.font.color.rgb = WHITE
            run.font.name = FONT
    return txBox

def add_subtitle_gray(slide, text, left=0.5, top=3.8):
    """Small gray text at bottom"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(9.0), Inches(0.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    run = tf.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(10)
    run.font.color.rgb = GRAY
    run.font.name = FONT

def add_picture_safe(slide, path, left, top, width=None, height=None):
    if os.path.exists(path):
        if width and height:
            slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(width), Inches(height))
        elif width:
            slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(width))
        else:
            slide.shapes.add_picture(path, Inches(left), Inches(top))


# ============================================================
# SLIDE 1 — HERO
# ============================================================
s = add_slide()
txBox = slide_shapes = s.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(6.5), Inches(2.5))
tf = txBox.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
run = p.add_run()
run.text = "Финансовый компас 1.0"
run.font.size = Pt(44)
run.font.bold = True
run.font.color.rgb = WHITE
run.font.name = FONT

p2 = tf.add_paragraph()
run2 = p2.add_run()
run2.text = "Пилотный практикум для предпринимателей\nна упрощённом налогообложении"
run2.font.size = Pt(18)
run2.font.color.rgb = GRAY
run2.font.name = FONT

add_subtitle_gray(s, "Онлайн · 10 встреч по 1-1,5 часа · 1-2 раза в неделю · Старт 6 апреля 2026", top=4.0)

# Stats in bottom area
stats_text = [
    ("10", "онлайн-встреч"),
    ("3", "эксперта"),
    ("1-1,5 ч", "длительность"),
    ("49 000 ₸", "стоимость"),
]
for i, (num, label) in enumerate(stats_text):
    x = 0.5 + i * 2.4
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(4.5), Inches(2.1), Inches(0.9))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x0D, 0x15, 0x2A)
    box.line.fill.background()

    tb = s.shapes.add_textbox(Inches(x), Inches(4.5), Inches(2.1), Inches(0.5))
    tf_s = tb.text_frame
    tf_s.word_wrap = True
    r = tf_s.paragraphs[0].add_run()
    r.text = num
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.color.rgb = BLUE
    r.font.name = FONT
    tf_s.paragraphs[0].alignment = PP_ALIGN.CENTER

    tb2 = s.shapes.add_textbox(Inches(x), Inches(4.95), Inches(2.1), Inches(0.4))
    tf_s2 = tb2.text_frame
    tf_s2.word_wrap = True
    r2 = tf_s2.paragraphs[0].add_run()
    r2.text = label
    r2.font.size = Pt(10)
    r2.font.color.rgb = GRAY
    r2.font.name = FONT
    tf_s2.paragraphs[0].alignment = PP_ALIGN.CENTER


# ============================================================
# SLIDE 2 — О СПИКЕРАХ: Малика Толессин
# ============================================================
s = add_slide()
add_title(s, "Малика Толессин, FCCA")
add_body(s, [
    ("Финансовый аналитик · FCCA · 20+ лет опыта", 16, True, WHITE),
    '',
    '• MSc Financial Economics — City University of London (2010-2011)',
    '• FCCA — Fellow Chartered Certified Accountant (полное членство ACCA)',
    '• PricewaterhouseCoopers — старший консультант, аудит по МСФО (2006-2010, Алматы)',
    '• Baker Hughes — управленческий учёт (2005-2006, Атырау)',
    '• Tengizchevroil — старший финансовый аналитик, затем тимлид по финансовой отчётности (2012-2018)',
    '• Chevron — FSS Upstream Analyst & Intercompany Team Lead (2018-2020)',
    '',
    '• Специализация: финансовая отчётность, управленческий учёт, коммерческий анализ, МСФО/IFRS',
], width=5.2, height=4.0)


# ============================================================
# SLIDE 3 — О СПИКЕРАХ: Гульнар Каримова
# ============================================================
s = add_slide()
add_title(s, "Гульнар Каримова")
add_body(s, [
    ("Главный бухгалтер · 30+ лет опыта", 16, True, WHITE),
    '',
    '• Магистр КИМЭП (1998)',
    '• Сертификат ACCA (2002)',
    '• Диплом профессионального бухгалтера (Италия, 2009)',
    '• USAID — финансовый координатор',
    '• Snamprogetti / ENI Group — бухгалтер, затем главный бухгалтер',
    '• KPO — старший бухгалтер',
    '• Agip / Baker Hughes — финансовые позиции',
    '• Финансовый координатор Каспийского региона (Казахстан, Кыргызстан, Узбекистан, Таджикистан)',
    '',
    '• Специализация: налогообложение ИП, упрощёнка 3%, форма 910, соцотчисления',
], width=5.2, height=4.0)


# ============================================================
# SLIDE 4 — О СПИКЕРАХ: Еркин Каримов
# ============================================================
s = add_slide()
add_title(s, "Еркин Каримов")
add_body(s, [
    ("Экономист · Создатель FinDirector · 20+ лет опыта", 16, True, WHITE),
    '',
    '• 2006-2012 — инженер-сметчик / QS в Agip KCO',
    '• 2014-2021 — ведущий инженер по оценке стоимости в LUKOIL',
    '• 2021-2023 — координатор по контролю проектов в TengizChevrOil',
    '• Контроль затрат на проектах стоимостью свыше US$2B (Chevron, LUKOIL, Agip/ENI, Kazakhmys)',
    '• С 2023 — предприниматель: ремонт квартир, общепит, продюсирование',
    '• Основатель креативной студии KARINA Media',
    '• 2024: разработал FinDirector — приложение для управления финансами малого бизнеса',
], width=5.2, height=4.0)

add_picture_safe(s, base + "101.png", 5.8, 1.3, width=4.2)


# ============================================================
# SLIDE 5 — БОЛИ
# ============================================================
s = add_slide()
add_title(s, "Почему предприниматели теряют деньги")
add_body(s, [
    ("80% малых бизнесов в Казахстане не ведут управленческий учёт", 14, True, GOLD),
    '',
    '• Кассовые разрывы — прибыль на бумаге есть, а платить зарплату и аренду нечем',
    '• Штрафы от налоговой — пропустили срок, ошибка в форме 910, штраф от 50 000 тг',
    '• Не знаете свою реальную прибыль — данные разбросаны по Excel, банку, приложению налоговой',
    '• Бизнес и личные деньги — одна куча. Себе зарплату не платите. Сколько зарабатывает бизнес — неизвестно',
    '',
    ("Результат: решения наугад, потеря денег и нервов", 14, True, RED),
])


# ============================================================
# SLIDE 6 — 5 ПОКАЗАТЕЛЕЙ
# ============================================================
s = add_slide()
add_title(s, "5 показателей для каждого владельца бизнеса")
add_body(s, [
    ("Если эти 5 инструментов работают — кассовых разрывов почти не бывает", 14, False, GRAY),
    '',
    '• 01. Остаток денег на счетах — сколько у вас есть прямо сейчас',
    '• 02. Денежный поток (Cash Flow) — куда реально уходят деньги',
    '• 03. Дебиторская задолженность — кто и сколько вам должен',
    '• 04. Платёжный календарь — кому и когда вы должны платить',
    '• 05. Прибыль (P&L) — сколько зарабатываете, тратите и сколько остаётся',
    '',
    ("На практикуме вы научитесь отслеживать каждый из них", 14, True, BLUE),
])


# ============================================================
# SLIDE 7 — ПРОГРАММА (встречи 01-03)
# ============================================================
s = add_slide()
add_title(s, "Программа: встречи 01-03")
add_body(s, [
    ("Финансовый фундамент", 20, True, WHITE),
    '',
    '• Разделение личных и бизнес-финансов',
    '• P&L из 5 строк — простейший отчёт о прибылях и убытках',
    '• Себестоимость vs операционные расходы',
    '• Простая структура отчёта для ИП и самозанятых',
    '',
    ("Формат: онлайн, 1-1,5 часа, 1-2 раза в неделю", 12, False, GRAY),
])


# ============================================================
# SLIDE 8 — ПРОГРАММА (встречи 04-06)
# ============================================================
s = add_slide()
add_title(s, "Программа: встречи 04-06")
add_body(s, [
    ("Учёт и налоги", 20, True, WHITE),
    '',
    '• Упрощёнка 3% для ИП и ТОО — как всё устроено',
    '• Обязательные соцплатежи и пенсионные отчисления',
    '• Платёжный календарь — когда и что платить',
    '• Как не попасть на штрафы от налоговой',
    '• Правильное заполнение формы 910',
])


# ============================================================
# SLIDE 9 — ПРОГРАММА (встречи 07-10)
# ============================================================
s = add_slide()
add_title(s, "Программа: встречи 07-10")
add_body(s, [
    ("Аналитика, рост и FinDirector на практике", 20, True, WHITE),
    '',
    '• Cash Flow и ДДС — полная картина движения денег',
    '• Юнит-экономика продукта — какие направления прибыльны',
    '• Снижение затрат — какие расходы инвестиции, а какие пустая трата',
    '• Работа с FinDirector: дашборд, договоры, ДДС, P&L',
    '• Настройка программы под ваш бизнес',
    '• Ваша обратная связь формирует продукт — вы первые, кто его использует',
])


# ============================================================
# SLIDE 10 — FINDIRECTOR: ДАШБОРД
# ============================================================
s = add_slide()
add_title(s, "FinDirector: Главный дашборд")
add_body(s, [
    '• Остатки на счетах, выручка, расходы и прибыль — всё на одном экране',
    '• Визуализация динамики выручки и расходов',
    '• AI-советы по управлению деньгами',
    '• Быстрый доступ ко всем разделам',
    '',
    ("Участники практикума получат 6 месяцев бесплатного доступа", 14, True, GREEN),
], width=5.0, height=3.5)

add_picture_safe(s, base + "101.png", 5.5, 1.0, width=4.5)


# ============================================================
# SLIDE 11 — FINDIRECTOR: ДДС
# ============================================================
s = add_slide()
add_title(s, "FinDirector: Движение денежных средств")
add_body(s, [
    '• Учёт всех платежей в одном окне',
    '• Категории: выручка, расходы, налоги, конвертации',
    '• Поддержка нескольких валют и статусов оплаты',
    '• Фильтры по периодам, категориям и контрагентам',
], width=5.0, height=3.0)

add_picture_safe(s, base + "103.png", 5.5, 1.0, width=4.5)


# ============================================================
# SLIDE 12 — FINDIRECTOR: P&L
# ============================================================
s = add_slide()
add_title(s, "FinDirector: Прибыли и убытки (P&L)")
add_body(s, [
    '• Анализ доходов, переменных и фиксированных расходов',
    '• Графики динамики прибыли по периодам',
    '• Выявление точек роста и слабых мест',
    '• Сводные таблицы по категориям расходов',
], width=5.0, height=3.0)

add_picture_safe(s, base + "104.png", 5.5, 1.0, width=4.5)


# ============================================================
# SLIDE 13 — СРАВНЕНИЕ ЦЕН
# ============================================================
s = add_slide()
add_title(s, "Стоимость ошибок vs стоимость практикума")

# Left column — market costs
add_body(s, [
    ("Без знаний и системы:", 16, True, RED),
    '',
    '• Штраф за ошибку в форме 910: от 50 000 ₸',
    '• Один кассовый разрыв: потеря клиента или просрочка',
    '• Бухгалтер на рынке: 5 000 - 15 000 ₸/час',
    '• Финансовый аналитик: 15 000 - 35 000 ₸/час',
    '• Даже 3-4 консультации = 60 000 - 120 000 ₸',
], left=0.5, top=1.3, width=4.5, height=3.5)

# Right column — our offer
add_body(s, [
    ("Пилотный практикум:", 16, True, GREEN),
    '',
    ("49 000 ₸", 28, True, GREEN),
    '',
    '• 10 онлайн-встреч по 1-1,5 часа',
    '• 3 эксперта: аналитик, бухгалтер, экономист',
    '• 6 месяцев бесплатного доступа к FinDirector',
    '• Эксклюзивная цена первого потока',
    '• Программа адаптируется лично под вас',
], left=5.3, top=1.3, width=4.5, height=3.5)


# ============================================================
# SLIDE 14 — ЧТО ПОЛУЧАТ УЧАСТНИКИ
# ============================================================
s = add_slide()
add_title(s, "Что вы получите как участник практикума")
add_body(s, [
    '• Базовые финансовые знания — P&L, Cash Flow, платёжный календарь, юнит-экономика',
    '• Защита от штрафов — за какими показателями следить, чтобы налоговая не стала проблемой',
    '• Контроль финансового здоровья — 5 ключевых показателей каждую неделю',
    '• 6 месяцев FinDirector бесплатно — программа адаптирована под участников практикума',
    '• Умение снижать затраты — какие расходы инвестиции, а какие пустая трата',
    '• Влияние на продукт — ваш голос напрямую влияет на разработку FinDirector',
])


# ============================================================
# SLIDE 15 — ПОЧЕМУ ПИЛОТНЫЙ ПРАКТИКУМ
# ============================================================
s = add_slide()
add_title(s, "Почему пилотный практикум — ваше преимущество")
add_body(s, [
    ("Вы — первый поток. Это значит:", 16, True, WHITE),
    '',
    '• Эксклюзивная цена 49 000 ₸ — в будущих потоках стоимость будет выше',
    '• Программа адаптируется лично под ваш бизнес и ваши вопросы',
    '• Вы формируете продукт FinDirector — он строится под реальные потребности',
    '• Максимальное внимание экспертов — малая группа, не поточная лекция',
    '• Прямой доступ к создателю программы и экспертам',
    '',
    ("FinDirector — новый продукт, и участники первого потока получают максимум внимания и влияния", 12, False, GRAY),
])


# ============================================================
# SLIDE 16 — CTA
# ============================================================
s = add_slide()
add_title(s, "Записаться в практикум", size=28)

txBox = s.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9.0), Inches(3.5))
tf = txBox.text_frame
tf.word_wrap = True

# Main headline
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Начните управлять финансами по цифрам, а не по ощущениям"
run.font.size = Pt(22)
run.font.bold = True
run.font.color.rgb = WHITE
run.font.name = FONT

# Price
p2 = tf.add_paragraph()
p2.space_before = Pt(16)
run = p2.add_run()
run.text = "49 000 ₸"
run.font.size = Pt(36)
run.font.bold = True
run.font.color.rgb = GREEN
run.font.name = FONT

# Details
details = [
    "10 онлайн-встреч по 1-1,5 часа · 1-2 раза в неделю",
    "3 эксперта: финансовый аналитик, бухгалтер, экономист",
    "6 месяцев бесплатного доступа к FinDirector",
    "Старт: 6 апреля 2026",
]
for d in details:
    p3 = tf.add_paragraph()
    p3.space_before = Pt(4)

    run_b = p3.add_run()
    run_b.text = "✓ "
    run_b.font.size = Pt(14)
    run_b.font.color.rgb = GREEN
    run_b.font.name = FONT

    run_t = p3.add_run()
    run_t.text = d
    run_t.font.size = Pt(14)
    run_t.font.color.rgb = WHITE
    run_t.font.name = FONT

# WhatsApp
p_wa = tf.add_paragraph()
p_wa.space_before = Pt(20)
run = p_wa.add_run()
run.text = "WhatsApp: +7 706 656 77 65"
run.font.size = Pt(18)
run.font.bold = True
run.font.color.rgb = BLUE
run.font.name = FONT

# Limited
add_subtitle_gray(s, "Пилотный практикум — количество мест ограничено", top=5.0)


# Save
output_path = base + "Финансовый_компас_1.0.pptx"
prs.save(output_path)
print(f"PPTX saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
